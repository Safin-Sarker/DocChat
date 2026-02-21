"""PPTX extraction utilities."""

import logging
from typing import Any, Dict, List

from pptx import Presentation

logger = logging.getLogger(__name__)


def _rows_to_markdown(rows: List[List[str]]) -> str:
    """Convert table rows to markdown text."""
    if not rows:
        return ""

    column_count = max(len(row) for row in rows)
    normalized_rows = [row + [""] * (column_count - len(row)) for row in rows]
    header = normalized_rows[0]
    body = normalized_rows[1:]

    header_line = "| " + " | ".join(cell or "" for cell in header) + " |"
    sep_line = "| " + " | ".join("---" for _ in header) + " |"
    body_lines = [
        "| " + " | ".join(cell or "" for cell in row) + " |"
        for row in body
    ]
    return "\n".join([header_line, sep_line] + body_lines)


class PptxExtractor:
    """Extract text and tables from PPTX files."""

    def extract_pages(self, pptx_path: str) -> List[Dict[str, Any]]:
        """Extract text from each slide as a page."""
        try:
            prs = Presentation(pptx_path)
        except Exception as exc:
            logger.error("PPTX load failed: %s", exc)
            return []

        pages: List[Dict[str, Any]] = []
        for slide_index, slide in enumerate(prs.slides):
            texts: List[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            slide_text = "\n".join(texts)
            if slide_text.strip():
                pages.append({"page_num": slide_index + 1, "text": slide_text})

        return pages

    def extract_tables(self, pptx_path: str) -> List[Dict[str, Any]]:
        """Extract tables from PPTX slides as markdown entries."""
        try:
            prs = Presentation(pptx_path)
        except Exception as exc:
            logger.error("PPTX load failed for tables: %s", exc)
            return []

        tables_out: List[Dict[str, Any]] = []
        table_counter = 0
        for slide_index, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    rows: List[List[str]] = []
                    for row in shape.table.rows:
                        rows.append([(cell.text or "").strip() for cell in row.cells])
                    markdown = _rows_to_markdown(rows)
                    if not markdown:
                        continue
                    tables_out.append({
                        "page": slide_index + 1,
                        "table_index": table_counter,
                        "markdown": markdown,
                    })
                    table_counter += 1

        return tables_out
